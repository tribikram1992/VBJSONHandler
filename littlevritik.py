from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

def add_slide(title, content, placeholder=True):
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    body.text = content
    if placeholder:
        left = Inches(5.5)
        top = Inches(2)
        width = Inches(3)
        height = Inches(3)
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height).text = "Image Placeholder"
    return slide

# Slide 1
slide_1 = prs.slides.add_slide(title_slide_layout)
slide_1.shapes.title.text = "All About Our Newborn"
slide_1.placeholders[1].text = "Vritik Deeptam Tripathy 🚀"

# Slides 2 to 11
add_slide("Welcome to Planet Baby!", 
          "📅 Launch Date: 5th May 2025\n🎉 Celebrated As: National Astronaut Day 🚀\n📍 Place of Birth: Cloudnine Hospital, Whitefield, Bangalore\n📏 Specs: Weight – 3.34 kg | Height – [To be filled]\n🕗 Time of Birth: 8:48:02 AM\n🎨 Complexion: Fair with pinkish palms & feet")

add_slide("Identity Crisis? Never Heard of It.",
          "Also Known As:\nYunay, Yunu, Gaja, Gajasana, Gajanan, Gajaraj, Mita, Kuni Baby, Vinayak, Lambodara, My Friend")

add_slide("First Press Conference",
          "🎤 First Cry: “Uaa… Uaa…”\n😊 Happy Noise: “Aaye”\n🍽️ Hungry Cry: “Aa…” with an oval-shaped face")

add_slide("Job Description - Baby Edition",
          "🔹 Professional Sleeper\n🔹 Part-time Crier\n🔹 Occasional Smiler\n🔹 Full-time Heart-Stealer 💘\n🔹 Hands: Tiny Boxer Fists 👊\n🔊 Hearing: Better than Siri!")

add_slide("Magical Month 1 – Our Little Wizard 🪄",
          "Nicknames: Pukul, Madhusudan, Dudulu, Sunulu, Munulu, Golu, Kuni Golu, Boss Baby\nLoves: Staring out of windows, lights, and walls")

add_slide("Vritik’s Top Passions",
          "💕 Eye Contact = Mind Control\n🎶 Loves Soothing Music\n👂 Recognizes Voices\n🪕 Dancing to Damburu!")

add_slide("Adventures of a 1-Month-Old Superhuman",
          "💤 Cutest Sleep Faces\n🚀 Blink-and-you-miss-it Movers\n🌀 Clockwise Twister\n👐 Fast Fist Pumps\n🧗 Lap-Climber\n🚶 Future Walker!")

add_slide("Baby’s Inner Circle",
          "👴 Aja & 👵 Aai – BFFs\n🎵 Lullabies: Tu Ku Musi, Tata Nadi")

add_slide("Mama’s Threats & Baby’s Reactions",
          "Mama: “Boarding school pathei debi!”\nBaby: Already packing my onesies!")

add_slide("Conclusion – Our Universe Has Shifted",
          "You've filled our lives with joy and wonder.\nWe can't wait for the magical mischief ahead!\n💖 Love you to the moon and back!")

# Save the presentation
prs.save("Vritik_Funny_Baby_Presentation.pptx")



